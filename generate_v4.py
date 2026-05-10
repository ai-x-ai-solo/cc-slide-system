#!/usr/bin/env python3
"""
V4: 社長設計の14枚固定構成スライド
1画面1感情 / 数字主役 / 余白活用 / セミナー登壇前提
"""

from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# カラーパレット（3色ルール）
BG       = RGBColor(0x0D, 0x0D, 0x0D)   # ほぼ黒
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)    # 白
PURPLE   = RGBColor(0x7B, 0x2F, 0xBE)   # パープル（アクセント）
PURPLE_L = RGBColor(0xA0, 0x5F, 0xD8)   # 薄パープル
GRAY     = RGBColor(0x88, 0x88, 0x88)   # グレー（ミュート）
DARK_CARD= RGBColor(0x1A, 0x1A, 0x2E)   # カード背景


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def rect(slide, l, t, w, h, color, line_color=None, line_w=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        if line_w:
            s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def t(slide, text, l, top, w, h, size, color,
      bold=False, align=PP_ALIGN.LEFT, wrap=True):
    if not text:
        return
    tb = slide.shapes.add_textbox(l, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Hiragino Kaku Gothic ProN"


def accent_lines(slide):
    """上下のアクセントライン"""
    rect(slide, 0, 0, SLIDE_W, Inches(0.06), PURPLE)
    rect(slide, 0, Inches(7.44), SLIDE_W, Inches(0.06), PURPLE)


# ===== 14枚のスライド =====

def slide01(prs):
    """AIを使える人と、使えない人の差がもう始まっています。"""
    s = blank(prs)
    bg(s)
    # パープルグラデーション風（左から右へ薄くなる帯）
    rect(s, 0, 0, Inches(5.0), SLIDE_H, PURPLE)
    rect(s, Inches(5.0), 0, Inches(4.0), SLIDE_H, RGBColor(0x4A, 0x1A, 0x7A))
    rect(s, Inches(9.0), 0, Inches(4.33), SLIDE_H, BG)
    # オーバーレイ（暗く）
    # 再度背景をセット
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    # パープル帯（左側）
    rect(s, 0, 0, Inches(0.3), SLIDE_H, PURPLE)
    rect(s, 0, 0, SLIDE_W, Inches(0.08), PURPLE)
    rect(s, 0, Inches(7.42), SLIDE_W, Inches(0.08), PURPLE)

    # メインコピー
    t(s, "AIを使える人と、", Inches(1.0), Inches(0.8),
      Inches(11.3), Inches(1.1), 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "使えない人の差が", Inches(1.0), Inches(1.8),
      Inches(11.3), Inches(1.1), 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "もう始まっています。", Inches(1.0), Inches(2.8),
      Inches(11.3), Inches(1.1), 42, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)

    # 区切り線
    rect(s, Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.04), GRAY)

    t(s, "月3,000円のAIツールで", Inches(1.0), Inches(4.3),
      Inches(11.3), Inches(0.7), 24, GRAY, align=PP_ALIGN.CENTER)
    t(s, "月10万円案件を狙える時代。", Inches(1.0), Inches(4.9),
      Inches(11.3), Inches(0.7), 24, WHITE, align=PP_ALIGN.CENTER)


def slide02(prs):
    """こんな作業、毎日やっていませんか？"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "こんな作業、毎日やっていませんか？", Inches(1.0), Inches(0.4),
      Inches(11.3), Inches(0.9), 32, WHITE, bold=True, align=PP_ALIGN.CENTER)

    items = ["資料作成", "SNS投稿", "リサーチ", "メール返信", "営業提案", "コピペ作業"]
    col_w = Inches(5.5)
    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        lx = Inches(1.5) + col * Inches(6.0)
        ty = Inches(1.5) + row * Inches(1.3)
        rect(s, lx, ty, col_w, Inches(1.0), DARK_CARD, PURPLE, 1)
        t(s, item, lx + Inches(0.3), ty + Inches(0.2),
          col_w - Inches(0.4), Inches(0.7), 26, WHITE, bold=True)

    rect(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.04), PURPLE)
    t(s, "「また同じことしてる…」", Inches(1.0), Inches(5.5),
      Inches(11.3), Inches(0.6), 22, GRAY, align=PP_ALIGN.CENTER)
    t(s, "その作業、AIで消せます。", Inches(1.0), Inches(6.1),
      Inches(11.3), Inches(0.7), 28, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)


def slide03(prs):
    """AIを"検索"にしか使えていない"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "でも多くの人は、", Inches(1.0), Inches(0.3),
      Inches(11.3), Inches(0.7), 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, 'AIを  "検索" にしか使えていない。', Inches(1.0), Inches(0.95),
      Inches(11.3), Inches(0.9), 34, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 左：ダメな例
    rect(s, Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.0), DARK_CARD, GRAY, 1)
    t(s, "😓 多くの人", Inches(0.7), Inches(2.1), Inches(5.1), Inches(0.6), 20, GRAY, bold=True)
    t(s, "ChatGPTに\n質問する", Inches(0.8), Inches(2.8),
      Inches(5.0), Inches(1.5), 30, GRAY, align=PP_ALIGN.CENTER)
    t(s, "→ 答えをもらうだけ", Inches(0.8), Inches(4.5),
      Inches(5.0), Inches(0.6), 18, GRAY, align=PP_ALIGN.CENTER)

    # 矢印
    t(s, "VS", Inches(6.1), Inches(3.5), Inches(1.1), Inches(0.8),
      36, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # 右：強い人
    rect(s, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.0), DARK_CARD, PURPLE, 2)
    t(s, "⚡ 強い人", Inches(7.5), Inches(2.1), Inches(5.1), Inches(0.6), 20, PURPLE_L, bold=True)
    t(s, "AIに\n働かせる", Inches(7.5), Inches(2.8),
      Inches(5.0), Inches(1.5), 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "→ 自動で仕事が回る", Inches(7.5), Inches(4.5),
      Inches(5.0), Inches(0.6), 18, WHITE, align=PP_ALIGN.CENTER)

    t(s, '本当に強い人は、AIに "働かせて" います。', Inches(1.0), Inches(6.3),
      Inches(11.3), Inches(0.7), 22, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)


def slide04(prs):
    """Claude Code ブートキャンプ 7日間"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "Claude Code ブートキャンプ", Inches(1.0), Inches(0.5),
      Inches(11.3), Inches(0.8), 30, GRAY, align=PP_ALIGN.CENTER)

    # 7 DAYS 超大
    t(s, "7", Inches(1.0), Inches(1.1), Inches(11.3), Inches(3.5),
      160, PURPLE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "DAYS", Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.2),
      64, WHITE, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(3.0), Inches(5.2), Inches(7.3), Inches(0.04), PURPLE)

    points = ["未経験OK", "日本語で操作", "実践型", "その場で完成"]
    for i, p in enumerate(points):
        lx = Inches(1.0) + i * Inches(3.0)
        t(s, f"✓ {p}", lx, Inches(5.4), Inches(2.8), Inches(0.7),
          20, WHITE, align=PP_ALIGN.CENTER)


def slide05(prs):
    """実際に学べる内容"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "実際に学べる内容", Inches(1.0), Inches(0.3),
      Inches(11.3), Inches(0.7), 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "AIエージェント ×", Inches(1.0), Inches(1.0),
      Inches(11.3), Inches(0.7), 26, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)

    left_items = ["LP制作", "案件探索", "スライド制作", "X投稿自動化"]
    right_items = ["動画編集", "LINE Bot構築", "システム開発", "SEO記事生成"]

    for i, item in enumerate(left_items):
        rect(s, Inches(0.5), Inches(1.9) + i * Inches(1.15),
             Inches(5.8), Inches(0.95), DARK_CARD, PURPLE, 1)
        t(s, item, Inches(0.8), Inches(2.0) + i * Inches(1.15),
          Inches(5.2), Inches(0.75), 24, WHITE, bold=True)

    for i, item in enumerate(right_items):
        rect(s, Inches(7.0), Inches(1.9) + i * Inches(1.15),
             Inches(5.8), Inches(0.95), DARK_CARD, PURPLE, 1)
        t(s, item, Inches(7.3), Inches(2.0) + i * Inches(1.15),
          Inches(5.2), Inches(0.75), 24, WHITE, bold=True)

    t(s, 'すべて 実務ベース。毎週テーマが増加中。', Inches(1.0), Inches(6.65),
      Inches(11.3), Inches(0.55), 18, GRAY, align=PP_ALIGN.CENTER)


def slide06(prs):
    """月4,500万円"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "「AI驚き屋」ではありません。", Inches(1.0), Inches(0.4),
      Inches(11.3), Inches(0.7), 26, GRAY, align=PP_ALIGN.CENTER)
    t(s, "実際にAIで", Inches(1.0), Inches(1.1),
      Inches(11.3), Inches(0.7), 28, WHITE, align=PP_ALIGN.CENTER)

    # 超巨大数字
    t(s, "月4,500万円", Inches(0.3), Inches(1.7),
      Inches(12.7), Inches(2.5), 80, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    t(s, "を生み出しているチームです。（4人で達成）", Inches(1.0), Inches(4.1),
      Inches(11.3), Inches(0.7), 24, WHITE, align=PP_ALIGN.CENTER)

    rect(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.04), PURPLE)

    proofs = ["100以上のシステム構築", "AIで講義300本 / 月 自動制作", "営業初月1,000万円案件獲得"]
    for i, p in enumerate(proofs):
        t(s, f"▶  {p}", Inches(2.5), Inches(5.1) + i * Inches(0.65),
          Inches(8.3), Inches(0.6), 22, GRAY)


def slide07(prs):
    """未経験でも大丈夫？"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "未経験でも大丈夫？", Inches(1.0), Inches(0.3),
      Inches(11.3), Inches(0.8), 32, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "大丈夫です。", Inches(1.0), Inches(1.1),
      Inches(11.3), Inches(1.0), 52, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # Before
    rect(s, Inches(0.4), Inches(2.3), Inches(5.5), Inches(3.5), DARK_CARD, GRAY, 1)
    t(s, "Before", Inches(0.6), Inches(2.4), Inches(5.0), Inches(0.6), 18, GRAY, bold=True)
    t(s, "「何から始めれば\nいいかわからない…」",
      Inches(0.6), Inches(3.0), Inches(5.1), Inches(1.8), 24, GRAY, align=PP_ALIGN.CENTER)

    t(s, "7日後 →", Inches(6.1), Inches(3.7), Inches(1.5), Inches(0.7),
      22, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # After
    rect(s, Inches(7.4), Inches(2.3), Inches(5.5), Inches(3.5), DARK_CARD, PURPLE, 2)
    t(s, "After", Inches(7.6), Inches(2.4), Inches(5.0), Inches(0.6), 18, PURPLE_L, bold=True)
    t(s, "「自分のシステムが\n動いてる！」",
      Inches(7.6), Inches(3.0), Inches(5.1), Inches(1.8), 24, WHITE, bold=True, align=PP_ALIGN.CENTER)

    t(s, "1期生：完全未経験 → 7日でシステム完成。", Inches(1.0), Inches(6.3),
      Inches(11.3), Inches(0.7), 22, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)


def slide08(prs):
    """朝「AIよくわからない」→夜「システム完成！」"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "ブートキャンプ当日", Inches(1.0), Inches(0.3),
      Inches(11.3), Inches(0.6), 24, GRAY, align=PP_ALIGN.CENTER)

    # 朝
    rect(s, Inches(0.5), Inches(1.1), Inches(5.5), Inches(2.2), DARK_CARD, GRAY, 1)
    t(s, "朝", Inches(0.7), Inches(1.2), Inches(5.0), Inches(0.5), 18, GRAY, bold=True)
    t(s, "「AIよくわからない…」", Inches(0.7), Inches(1.8),
      Inches(5.1), Inches(1.2), 26, GRAY, align=PP_ALIGN.CENTER)

    # 矢印
    t(s, "↓", Inches(0.5), Inches(3.4), Inches(5.5), Inches(0.8),
      40, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # 夜
    rect(s, Inches(0.5), Inches(4.2), Inches(5.5), Inches(2.3), DARK_CARD, PURPLE, 2)
    t(s, "夜", Inches(0.7), Inches(4.3), Inches(5.0), Inches(0.5), 18, PURPLE_L, bold=True)
    t(s, "「自分のシステム\n完成した！」", Inches(0.7), Inches(4.8),
      Inches(5.1), Inches(1.5), 30, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 右側サポート
    supports = ["リアルタイム作業", "その場で質問可能", "つまずきを即解決"]
    for i, sp in enumerate(supports):
        rect(s, Inches(7.0), Inches(1.5) + i * Inches(1.4),
             Inches(5.8), Inches(1.1), DARK_CARD, PURPLE, 1)
        t(s, f"✓  {sp}", Inches(7.3), Inches(1.65) + i * Inches(1.4),
          Inches(5.3), Inches(0.8), 24, WHITE)

    t(s, "「終わらなかった」が起きにくい構造。", Inches(1.0), Inches(6.6),
      Inches(11.3), Inches(0.6), 20, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)


def slide09(prs):
    """あなたが受け取るもの 合計2,485,000円相当"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "あなたが受け取るもの", Inches(1.0), Inches(0.3),
      Inches(11.3), Inches(0.7), 28, WHITE, bold=True, align=PP_ALIGN.CENTER)

    items_left = ["講義13テーマ（毎週増加中）", "専用コミュニティ", "ブートキャンプ4回"]
    items_right = ["社内システム12種", "1on1コンサル（2時間）", "補講サポート1ヶ月"]

    for i, item in enumerate(items_left):
        t(s, f"✓  {item}", Inches(0.8), Inches(1.2) + i * Inches(0.75),
          Inches(5.8), Inches(0.65), 20, GRAY)
    for i, item in enumerate(items_right):
        t(s, f"✓  {item}", Inches(7.0), Inches(1.2) + i * Inches(0.75),
          Inches(5.8), Inches(0.65), 20, GRAY)

    rect(s, Inches(1.0), Inches(3.7), Inches(11.3), Inches(0.05), PURPLE)
    t(s, "合計", Inches(1.0), Inches(3.9), Inches(11.3), Inches(0.65),
      26, GRAY, align=PP_ALIGN.CENTER)
    t(s, "¥2,485,000相当", Inches(0.3), Inches(4.5),
      Inches(12.7), Inches(2.0), 72, PURPLE, bold=True, align=PP_ALIGN.CENTER)


def slide10(prs):
    """しかし今回 月々約7,532円〜"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "しかし今回、2期生限定で。", Inches(1.0), Inches(0.5),
      Inches(11.3), Inches(0.7), 26, GRAY, align=PP_ALIGN.CENTER)
    t(s, "クレジット分割なら", Inches(1.0), Inches(1.4),
      Inches(11.3), Inches(0.7), 26, WHITE, align=PP_ALIGN.CENTER)

    # 月々価格
    t(s, "月々 約7,532円〜", Inches(0.3), Inches(2.1),
      Inches(12.7), Inches(2.2), 72, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(3.0), Inches(4.4), Inches(7.3), Inches(0.04), GRAY)

    t(s, "1日あたり 約251円", Inches(1.0), Inches(4.6),
      Inches(11.3), Inches(0.8), 32, WHITE, align=PP_ALIGN.CENTER)

    t(s, "コーヒー1杯以下で、月10万円案件を狙えるスキルが身につく。", Inches(1.0), Inches(5.7),
      Inches(11.3), Inches(0.7), 20, GRAY, align=PP_ALIGN.CENTER)


def slide11(prs):
    """198,000円。高いと思いましたか？"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "198,000円。", Inches(1.0), Inches(0.4),
      Inches(11.3), Inches(1.1), 52, WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "高いと思いましたか？", Inches(1.0), Inches(1.5),
      Inches(11.3), Inches(0.7), 28, GRAY, align=PP_ALIGN.CENTER)

    t(s, "でももし、", Inches(1.0), Inches(2.5),
      Inches(11.3), Inches(0.6), 22, WHITE, align=PP_ALIGN.CENTER)

    ifs = ["月10万円案件を2件取ったら", "AIで業務を週10時間削減できたら", "自動化スキルで副業が軌道に乗ったら"]
    for i, line in enumerate(ifs):
        t(s, f"▶  {line}", Inches(2.5), Inches(3.2) + i * Inches(0.75),
          Inches(8.3), Inches(0.65), 22, WHITE)

    rect(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.04), PURPLE)
    t(s, "回収できる金額は、それ以上かもしれません。", Inches(1.0), Inches(5.7),
      Inches(11.3), Inches(0.7), 26, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)


def slide12(prs):
    """推定0.1%未満"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "Claude Codeを仕事に使いこなしている日本人は", Inches(0.5), Inches(0.4),
      Inches(12.3), Inches(0.7), 24, GRAY, align=PP_ALIGN.CENTER)

    # 超巨大0.1%
    t(s, "0.1%", Inches(0.3), Inches(0.9),
      Inches(12.7), Inches(3.8), 148, PURPLE, bold=True, align=PP_ALIGN.CENTER)

    t(s, "未満", Inches(0.3), Inches(4.5),
      Inches(12.7), Inches(0.9), 48, WHITE, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.05), PURPLE)

    t(s, "今はまだ、「早い者勝ち」の段階です。", Inches(1.0), Inches(5.7),
      Inches(11.3), Inches(0.7), 26, PURPLE_L, bold=True, align=PP_ALIGN.CENTER)
    t(s, "先に動いた人だけが、圧倒的に有利になる。", Inches(1.0), Inches(6.4),
      Inches(11.3), Inches(0.7), 22, GRAY, align=PP_ALIGN.CENTER)


def slide13(prs):
    """1年後、どちら側にいますか？"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "もし1年後、AIを使える人だけが生産性を10倍にしていたら？", Inches(0.5), Inches(0.3),
      Inches(12.3), Inches(0.8), 22, GRAY, align=PP_ALIGN.CENTER)

    # 左：使えなかった側
    rect(s, Inches(0.4), Inches(1.3), Inches(5.5), Inches(4.5), DARK_CARD, GRAY, 1)
    t(s, "そのままの自分", Inches(0.6), Inches(1.45), Inches(5.0), Inches(0.6), 20, GRAY, bold=True)
    futures_l = ["今と同じ作業を繰り返す", "AIが得意な人に仕事を奪われる", "差が広がるのを見ている"]
    for i, f in enumerate(futures_l):
        t(s, f"✗  {f}", Inches(0.7), Inches(2.2) + i * Inches(0.9),
          Inches(5.0), Inches(0.8), 20, GRAY)

    t(s, "あなたは\nどちら側\nにいますか？", Inches(6.0), Inches(2.8),
      Inches(1.5), Inches(2.5), 20, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 右：使える側
    rect(s, Inches(7.5), Inches(1.3), Inches(5.5), Inches(4.5), DARK_CARD, PURPLE, 2)
    t(s, "AIを使いこなす自分", Inches(7.7), Inches(1.45), Inches(5.0), Inches(0.6), 20, PURPLE_L, bold=True)
    futures_r = ["自動化で時間が生まれる", "月10万円以上の案件を受注", "スキルで自由な働き方を実現"]
    for i, f in enumerate(futures_r):
        t(s, f"✓  {f}", Inches(7.8), Inches(2.2) + i * Inches(0.9),
          Inches(5.0), Inches(0.8), 20, WHITE)

    t(s, "決断するなら、「やってみたい」と思っている今です。", Inches(1.0), Inches(6.3),
      Inches(11.3), Inches(0.7), 20, GRAY, align=PP_ALIGN.CENTER)


def slide14(prs):
    """CTA"""
    s = blank(prs)
    bg(s)
    accent_lines(s)

    t(s, "本質のClaude Code 完全攻略", Inches(1.0), Inches(0.8),
      Inches(11.3), Inches(0.7), 26, GRAY, align=PP_ALIGN.CENTER)
    t(s, "7dayブートキャンプ", Inches(1.0), Inches(1.5),
      Inches(11.3), Inches(1.0), 44, WHITE, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.04), PURPLE)

    t(s, "未経験から、AIで仕事を作れる人へ。", Inches(1.0), Inches(2.9),
      Inches(11.3), Inches(0.7), 24, GRAY, align=PP_ALIGN.CENTER)

    # CTAボタン
    rect(s, Inches(2.5), Inches(4.0), Inches(8.3), Inches(1.4), PURPLE)
    t(s, "本日参加受付中 →", Inches(2.7), Inches(4.2),
      Inches(7.9), Inches(1.0), 36, WHITE, bold=True, align=PP_ALIGN.CENTER)

    t(s, "198,000円（税抜）　クレジット分割：月々約7,532円〜", Inches(1.0), Inches(5.7),
      Inches(11.3), Inches(0.6), 20, GRAY, align=PP_ALIGN.CENTER)
    t(s, "※ 2期生限定価格", Inches(1.0), Inches(6.35),
      Inches(11.3), Inches(0.5), 16, GRAY, align=PP_ALIGN.CENTER)


def main():
    prs = new_prs()
    slides = [
        slide01, slide02, slide03, slide04, slide05, slide06, slide07,
        slide08, slide09, slide10, slide11, slide12, slide13, slide14,
    ]
    for fn in slides:
        fn(prs)

    out_dir = Path(__file__).parent / "output"
    out_dir.mkdir(exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = str(out_dir / f"cc_v4_structured_{ts}.pptx")
    prs.save(out)
    print(f"✅ V4生成完了: {out}")
    print(f"   14枚 / ダーク×パープル / 1スライド1感情")


if __name__ == "__main__":
    main()
