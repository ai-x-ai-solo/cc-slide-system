#!/usr/bin/env python3
"""
V6: 3ターゲット訴求パターン
  Pattern A: 副業・案件獲得訴求（黒×パープル）
  Pattern B: 法人・業務改善訴求（白×青）
  Pattern C: 初心者安心訴求（白×グリーン）
設計: 1スライド1メッセージ / 1画面1感情 / 感情導線で構成
"""

from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
ASSETS_DIR      = Path(__file__).parent / "assets" / "references"
PROCESSED_DIR   = Path(__file__).parent / "assets" / "processed"

# ===== 3ターゲット訴求パターン =====
PATTERNS = {
    "A": {
        "name": "副業・案件獲得訴求型（黒×パープル）",
        "target": "AI副業に興味がある個人・スキルプラス受講生",
        "emotion": "今動けば先行者になれる",
        "bg":       RGBColor(0x0D, 0x0D, 0x0D),
        "primary":  RGBColor(0x7B, 0x2F, 0xFF),   # パープル
        "accent":   RGBColor(0xAA, 0x66, 0xFF),   # 淡パープル
        "text":     RGBColor(0xFF, 0xFF, 0xFF),
        "text_inv": RGBColor(0x0D, 0x0D, 0x0D),
        "muted":    RGBColor(0x88, 0x88, 0x99),
        "card":     RGBColor(0x1A, 0x0D, 0x2E),
        "card_bd":  RGBColor(0x7B, 0x2F, 0xFF),
        "safe":     RGBColor(0x44, 0xFF, 0x88),
        "danger":   RGBColor(0xFF, 0x44, 0x44),
    },
    "B": {
        "name": "法人・業務改善訴求型（白×青）",
        "target": "中小企業経営者・管理職・業務改善担当者",
        "emotion": "自社の業務に使えそう",
        "bg":       RGBColor(0xFF, 0xFF, 0xFF),
        "primary":  RGBColor(0x1A, 0x5F, 0xBF),   # 濃い青
        "accent":   RGBColor(0x3A, 0x9B, 0xFF),   # 明るい青
        "text":     RGBColor(0x1A, 0x1A, 0x2E),
        "text_inv": RGBColor(0xFF, 0xFF, 0xFF),
        "muted":    RGBColor(0x66, 0x77, 0x88),
        "card":     RGBColor(0xF0, 0xF5, 0xFF),
        "card_bd":  RGBColor(0xB0, 0xCC, 0xFF),
        "safe":     RGBColor(0x1A, 0x8A, 0x3A),
        "danger":   RGBColor(0xCC, 0x22, 0x22),
    },
    "C": {
        "name": "初心者安心訴求型（白×グリーン）",
        "target": "AI初心者・PC初心者・未経験者",
        "emotion": "自分でもできそう",
        "bg":       RGBColor(0xFA, 0xFF, 0xFA),
        "primary":  RGBColor(0x1A, 0xA0, 0x6E),   # グリーン
        "accent":   RGBColor(0xFF, 0x7A, 0x45),   # コーラル（温かみ）
        "text":     RGBColor(0x1A, 0x2E, 0x1A),
        "text_inv": RGBColor(0xFF, 0xFF, 0xFF),
        "muted":    RGBColor(0x77, 0x99, 0x77),
        "card":     RGBColor(0xE8, 0xF8, 0xF0),
        "card_bd":  RGBColor(0x88, 0xDD, 0xAA),
        "safe":     RGBColor(0x1A, 0xA0, 0x6E),
        "danger":   RGBColor(0xFF, 0x7A, 0x45),
    },
}


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, p):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = p["bg"]


def box(slide, l, t, w, h, color, bd_color=None, bd_w=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if bd_color:
        s.line.color.rgb = bd_color
        if bd_w:
            s.line.width = Pt(bd_w)
    else:
        s.line.fill.background()
    return s


def tx(slide, text, l, top, w, h, size, color,
       bold=False, align=PP_ALIGN.LEFT, wrap=True):
    if not text:
        return
    tb = slide.shapes.add_textbox(l, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    r = para.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Hiragino Kaku Gothic ProN"


def frame(slide, p):
    box(slide, 0, 0, SLIDE_W, Inches(0.07), p["primary"])
    box(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), p["accent"])


def add_photo(slide, filename, left, top, width, height):
    """正規化済み画像を優先、なければ元画像、なければスキップ"""
    normalized = filename.replace(".png", "-normalized.png")
    proc = PROCESSED_DIR / "profiles" / normalized
    ref  = ASSETS_DIR   / "profiles" / filename
    path = proc if proc.exists() else (ref if ref.exists() else None)
    if path:
        slide.shapes.add_picture(str(path), left, top, width, height)
        return True
    return False


def add_actionmap(slide, filename, left, top, width, height):
    """正規化済みアクションマップ画像を挿入"""
    normalized = filename.replace(".png", "-normalized.png")
    proc = PROCESSED_DIR / "actionmap" / normalized
    ref  = ASSETS_DIR   / "actionmap" / filename
    path = proc if proc.exists() else (ref if ref.exists() else None)
    if path:
        slide.shapes.add_picture(str(path), left, top, width, height)
        return True
    return False


# ===== スライド関数（Pattern A: 副業・案件獲得） =====

def sA01_cover(slide, p):
    bg(slide, p)
    box(slide, 0, 0, SLIDE_W, Inches(0.07), p["primary"])
    box(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), p["accent"])
    tx(slide, "7日で、", Inches(1.0), Inches(1.0), Inches(11.3), Inches(1.2),
       52, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "月10万円案件を狙えるスキルへ。", Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.8),
       56, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(2.0), Inches(4.35), Inches(9.3), Inches(0.04), p["accent"])
    tx(slide, "本質のClaude Code 完全攻略 7dayブートキャンプ",
       Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.7), 22, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "アドネス株式会社", Inches(1.0), Inches(6.8),
       Inches(11.3), Inches(0.5), 16, p["muted"], align=PP_ALIGN.RIGHT)


def sA02_problem(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "毎日同じ作業を繰り返していませんか？",
       Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8), 26, p["muted"],
       align=PP_ALIGN.CENTER)
    items = [
        "資料作成に毎回2〜3時間かかっている",
        "同じメールを少しだけ変えて何度も送っている",
        "副業を始めたいが、何のスキルを身につければいいか分からない",
        "「これ、誰でもできる作業だよな」と感じながら続けている",
    ]
    for i, item in enumerate(items):
        ty = Inches(1.4) + i * Inches(1.3)
        box(slide, Inches(0.8), ty, Inches(0.5), Inches(0.9), p["primary"])
        tx(slide, str(i + 1), Inches(0.82), ty + Inches(0.1),
           Inches(0.45), Inches(0.75), 28, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        tx(slide, item, Inches(1.5), ty + Inches(0.12),
           Inches(11.0), Inches(0.75), 22, p["text"])


def sA03_gap(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "AI格差は、今この瞬間も広がっています。",
       Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.75), 26, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(0.5), Inches(1.25), Inches(5.8), Inches(5.0), p["card"], p["card_bd"], 1)
    tx(slide, "AIを使えない人", Inches(0.7), Inches(1.4), Inches(5.3), Inches(0.6),
       18, p["muted"], bold=True)
    for i, t in enumerate(["手作業・繰り返し作業", "月収は横ばい", "スキルが身につかない"]):
        tx(slide, f"✗  {t}", Inches(0.9), Inches(2.1) + i * Inches(0.9),
           Inches(5.1), Inches(0.8), 20, p["muted"])
    tx(slide, "VS", Inches(6.45), Inches(3.2), Inches(0.9), Inches(1.0),
       32, p["accent"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(7.0), Inches(1.25), Inches(5.8), Inches(5.0), p["card"], p["primary"], 2)
    tx(slide, "AIを使いこなす人", Inches(7.2), Inches(1.4), Inches(5.3), Inches(0.6),
       18, p["primary"], bold=True)
    for i, t in enumerate(["作業を自動化して時間を生む", "月10万円以上の案件を獲得", "スキルで選ばれる側になる"]):
        tx(slide, f"✓  {t}", Inches(7.4), Inches(2.1) + i * Inches(0.9),
           Inches(5.1), Inches(0.8), 20, p["text"])
    tx(slide, "→ 差がつくのは「知っているかどうか」だけ。",
       Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.65), 18, p["accent"],
       align=PP_ALIGN.CENTER)


def sA04_solution(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Codeを使いこなしている日本人は",
       Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.65), 22, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "推定", Inches(0.8), Inches(1.1), Inches(3.5), Inches(1.5),
       36, p["muted"], align=PP_ALIGN.RIGHT)
    tx(slide, "0.1%", Inches(4.0), Inches(0.8), Inches(5.5), Inches(2.2),
       120, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "未満", Inches(9.5), Inches(1.1), Inches(3.0), Inches(1.5),
       36, p["muted"])
    box(slide, Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.04), p["accent"])
    tx(slide, "今から始めれば、あなたは「使いこなせる側」になれます。",
       Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.75), 24, p["text"],
       align=PP_ALIGN.CENTER)
    tx(slide, "Claude Codeは日本語で動きます。プログラミングは不要。",
       Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.7), 22, p["accent"],
       bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "→ 早い者勝ちのスキル。今が、最後のタイミングかもしれない。",
       Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.65), 17, p["muted"],
       align=PP_ALIGN.CENTER)


def sA05_mikami(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR 01", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "みかみ", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　代表取締役", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    points = [
        ("東大在学中に起業 → AIで組織を最適化", "規模"),
        ("AIで月商5億円を達成", "実績"),
        ("スキルプラス受講生 6,000名以上", "信頼"),
        ("「実際の事業でAIを使い続けている側」", ""),
    ]
    for i, (point, tag) in enumerate(points):
        ty = Inches(2.65) + i * Inches(1.0)
        if tag:
            box(slide, Inches(0.8), ty, Inches(1.1), Inches(0.7), p["primary"])
            tx(slide, tag, Inches(0.82), ty + Inches(0.1), Inches(1.05), Inches(0.55),
               13, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        tx(slide, point, Inches(2.1), ty + Inches(0.1), Inches(6.7), Inches(0.7),
           22, p["text"])
    add_photo(slide, "mikami-profile.png", Inches(10.1), Inches(0.55), Inches(3.0), Inches(3.0))


def sA06_sato(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR 02", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "佐藤将司", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　AI事業責任者", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    story = [
        "元・公務員（月給12万円）／プログラミング経験ゼロ",
        "↓　Claude Codeで100以上のシステムを構築",
        "↓　営業初月　1,000万円案件を獲得",
        "↓　AIで月4,500万円規模の事業を運営",
    ]
    colors = [p["muted"], p["text"], p["text"], p["primary"]]
    for i, (line, color) in enumerate(zip(story, colors)):
        tx(slide, line, Inches(1.0), Inches(2.7) + i * Inches(1.0),
           Inches(8.0), Inches(0.85), 22 if i < 3 else 24, color, bold=(i == 3))
    add_photo(slide, "sato-profile.png", Inches(10.1), Inches(0.55), Inches(3.0), Inches(3.0))


def sA07_bootcamp(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Code ブートキャンプ", Inches(1.0), Inches(0.3),
       Inches(11.3), Inches(0.65), 26, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "7", Inches(1.0), Inches(0.95), Inches(11.3), Inches(3.5),
       160, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "DAYS", Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.1),
       64, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(0.05), p["accent"])
    for i, pt in enumerate(["未経験OK", "日本語で操作", "実践型", "その場で完成"]):
        tx(slide, f"✓ {pt}", Inches(0.8) + i * Inches(3.0), Inches(5.4),
           Inches(2.8), Inches(0.7), 20, p["text"], align=PP_ALIGN.CENTER)


def sA08_actionmap(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "受講後、こんなことができるようになります",
       Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.65), 24, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    themes = [
        ("actionmap-kensakujidoka.png",  "案件探索\n自動化", "営業前に\n案件候補が\n自動で集まる"),
        ("actionmap-slide-seisaku.png",  "スライド\n制作",   "提案資料を\n半自動で\n作れる"),
        ("actionmap-system-kaihatsu.png","システム\n開発",   "業務改善ツールを\nAIで\n自分で構築"),
    ]
    col_w = Inches(4.0)
    img_h = Inches(2.25)
    for i, (fname, theme, value) in enumerate(themes):
        lx = Inches(0.55) + i * Inches(4.27)
        # 画像サムネイル
        has_img = add_actionmap(slide, fname, lx, Inches(1.0), col_w, img_h)
        if not has_img:
            box(slide, lx, Inches(1.0), col_w, img_h, p["card"], p["card_bd"], 1)
            tx(slide, theme, lx, Inches(1.5), col_w, Inches(1.3),
               22, p["primary"], bold=True, align=PP_ALIGN.CENTER)
        # テーマ名
        box(slide, lx, Inches(3.35), col_w, Inches(0.6), p["primary"])
        tx(slide, theme.replace("\n", " "), lx + Inches(0.1), Inches(3.38),
           col_w - Inches(0.2), Inches(0.55), 16, p["text_inv"], bold=True)
        # 翻訳価値
        box(slide, lx, Inches(4.0), col_w, Inches(2.2), p["card"], p["card_bd"], 1)
        tx(slide, value, lx + Inches(0.15), Inches(4.15),
           col_w - Inches(0.3), Inches(1.9), 22, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, "全13テーマ・毎週増加中", Inches(0.8), Inches(6.85),
       Inches(11.7), Inches(0.45), 15, p["muted"], align=PP_ALIGN.RIGHT)


def sA09_proof(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "1期生の実績", Inches(0.8), Inches(0.2),
       Inches(11.7), Inches(0.7), 28, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(4.0), Inches(1.05), Inches(5.3), Inches(0.05), p["accent"])
    facts = [
        ("完全未経験から", "7日", "でシステムを完成"),
        ("LP制作・SNS自動化・業務効率化", "各自が完成", ""),
        ("コミュニティで助け合いながら", "スキルアップ中", ""),
    ]
    for i, (pre, num, post) in enumerate(facts):
        ty = Inches(1.3) + i * Inches(1.7)
        tx(slide, pre, Inches(0.8), ty, Inches(11.7), Inches(0.6), 20, p["muted"],
           align=PP_ALIGN.CENTER)
        tx(slide, num, Inches(0.8), ty + Inches(0.55), Inches(11.7), Inches(0.9),
           42, p["primary"], bold=True, align=PP_ALIGN.CENTER)
        if post:
            tx(slide, post, Inches(0.8), ty + Inches(1.4), Inches(11.7), Inches(0.5),
               18, p["muted"], align=PP_ALIGN.CENTER)


def sA10_value(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "受け取れる価値", Inches(0.8), Inches(0.2),
       Inches(11.7), Inches(0.65), 24, p["muted"], align=PP_ALIGN.CENTER)
    items = [
        ("講義（全13テーマ・動画+テキスト）", "349,000円"),
        ("専用コミュニティ＋質問＋アーカイブ", "100,000円"),
        ("ブートキャンプ（全4回・各2時間）", "200,000円"),
        ("社内実用システム12種（そのまま使える）", "超高額"),
    ]
    for i, (item, val) in enumerate(items):
        ty = Inches(1.05) + i * Inches(1.0)
        box(slide, Inches(0.6), ty, Inches(8.8), Inches(0.85), p["card"], p["card_bd"], 1)
        tx(slide, item, Inches(0.8), ty + Inches(0.12), Inches(7.5), Inches(0.65), 19, p["text"])
        box(slide, Inches(9.6), ty, Inches(3.1), Inches(0.85), p["primary"])
        tx(slide, val, Inches(9.65), ty + Inches(0.1), Inches(3.0), Inches(0.65),
           17, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.05), p["accent"])
    tx(slide, "総額", Inches(0.6), Inches(5.3), Inches(4.0), Inches(0.9), 24, p["muted"],
       align=PP_ALIGN.RIGHT)
    tx(slide, "2,485,000", Inches(4.7), Inches(5.2), Inches(6.5), Inches(1.1),
       60, p["primary"], bold=True)
    tx(slide, "円相当", Inches(11.3), Inches(5.5), Inches(1.8), Inches(0.7), 22, p["muted"])


def sA11_price(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "価格", Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.65),
       22, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "198,000", Inches(0.8), Inches(1.0), Inches(11.7), Inches(2.5),
       110, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "円（税抜）", Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.75),
       30, p["muted"], align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.05), p["primary"])
    tx(slide, "1日あたり", Inches(1.5), Inches(4.4), Inches(4.5), Inches(0.65),
       22, p["muted"], align=PP_ALIGN.RIGHT)
    tx(slide, "251円", Inches(6.0), Inches(4.3), Inches(5.0), Inches(0.85),
       42, p["accent"], bold=True)
    tx(slide, "クレジット分割：月々約 7,532円〜", Inches(0.8), Inches(5.4),
       Inches(11.7), Inches(0.65), 22, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "2期生限定価格", Inches(0.8), Inches(6.2),
       Inches(11.7), Inches(0.5), 18, p["primary"], bold=True, align=PP_ALIGN.CENTER)


def sA12_roi(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "投資対効果", Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.65),
       22, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "月10万円の案件を", Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8),
       30, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "2件", Inches(0.8), Inches(1.85), Inches(11.7), Inches(1.8),
       90, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "取るだけで全額回収", Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.8),
       30, p["muted"], align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "法人AIコンサルの相場：月50〜200万円",
       Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.65), 20, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "月10万円案件を狙うための実践スキルと制作経験を、ここで積み上げます。",
       Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.65), 20, p["text"],
       align=PP_ALIGN.CENTER)


def sA13_urgency(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Codeを使いこなしている日本人", Inches(0.8), Inches(0.6),
       Inches(11.7), Inches(0.65), 22, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "推定 0.1%未満", Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.5),
       64, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "先に動いた人だけが、圧倒的に有利になる。",
       Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.9), 32, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "AIスキルは「知っているかどうか」の差です。今この瞬間が、最も早いタイミングです。",
       Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8), 20, p["muted"],
       align=PP_ALIGN.CENTER)


def sA14_cta(slide, p):
    bg(slide, p)
    box(slide, 0, 0, SLIDE_W, Inches(0.07), p["primary"])
    box(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), p["accent"])
    tx(slide, "7日で、月10万円案件を狙えるスキルへ。",
       Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.9), 26, p["muted"],
       align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "本質のClaude Code 完全攻略", Inches(0.8), Inches(2.1),
       Inches(11.7), Inches(0.75), 30, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, "7dayブートキャンプ", Inches(0.8), Inches(2.8),
       Inches(11.7), Inches(1.2), 52, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.5), p["primary"])
    tx(slide, "今すぐ申し込む →", Inches(2.0), Inches(4.63),
       Inches(9.3), Inches(0.65), 38, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "198,000円（税抜）　月々 7,532円〜　2期生限定価格",
       Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6), 18, p["muted"],
       align=PP_ALIGN.CENTER)


# ===== スライドリスト（Pattern A） =====
SLIDES_A = [
    sA01_cover, sA02_problem, sA03_gap, sA04_solution,
    sA05_mikami, sA06_sato, sA07_bootcamp, sA08_actionmap,
    sA09_proof, sA10_value, sA11_price, sA12_roi,
    sA13_urgency, sA14_cta,
]


# ===== Pattern B: 法人・業務改善訴求 =====

def sB01_cover(slide, p):
    bg(slide, p)
    frame(slide, p)
    box(slide, 0, Inches(2.8), SLIDE_W, Inches(2.2), p["card"])
    tx(slide, "AI導入に、", Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.9),
       34, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "外注は必要ない。", Inches(1.0), Inches(1.85), Inches(11.3), Inches(1.2),
       52, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "本質のClaude Code 完全攻略 7dayブートキャンプ",
       Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.7), 22, p["primary"],
       bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "7日間で、自社のAI活用を内製化する。",
       Inches(1.0), Inches(3.9), Inches(11.3), Inches(0.65), 20, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "アドネス株式会社", Inches(1.0), Inches(6.8),
       Inches(11.3), Inches(0.5), 16, p["muted"], align=PP_ALIGN.RIGHT)


def sB02_problem(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "社内の繰り返し作業に、毎月いくらかかっていますか？",
       Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.75), 24, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("資料作成・メール対応・データ整理", "月30〜50時間の工数"),
        ("AI導入を検討しているが何から始めればいいか分からない", "機会損失"),
        ("外注すると高い・社内に知識がない・誰に頼めばいいか\n分からない", "コスト増大"),
    ]
    for i, (prob, tag) in enumerate(items):
        ty = Inches(1.3) + i * Inches(1.6)
        box(slide, Inches(0.6), ty, Inches(9.0), Inches(1.3), p["card"], p["card_bd"], 1)
        tx(slide, prob, Inches(0.8), ty + Inches(0.15), Inches(8.5), Inches(0.8), 20, p["text"])
        box(slide, Inches(9.8), ty, Inches(3.0), Inches(1.3), p["primary"])
        tx(slide, tag, Inches(9.8), ty + Inches(0.25), Inches(3.0), Inches(0.8),
           15, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "「AIを入れたい」と思っているだけで、動けていない企業が多い。",
       Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.65), 18, p["muted"],
       align=PP_ALIGN.CENTER)


def sB03_solution(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Codeとは", Inches(0.8), Inches(0.2),
       Inches(11.7), Inches(0.65), 22, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "日本語で動かせるAI開発ツール",
       Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.1),
       44, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(2.25), Inches(10.3), Inches(0.05), p["accent"])
    feats = [
        ("プログラミング不要", "日本語で指示するだけで動く"),
        ("業務直結", "実際の社内システムを\nそのまま構築"),
        ("スピード", "繰り返し作業を、24時間動くAIに任せられる"),
    ]
    for i, (title, desc) in enumerate(feats):
        lx = Inches(0.6) + i * Inches(4.2)
        box(slide, lx, Inches(2.5), Inches(3.9), Inches(3.5), p["card"], p["card_bd"], 1)
        tx(slide, title, lx + Inches(0.15), Inches(2.7), Inches(3.6), Inches(0.7),
           20, p["primary"], bold=True, align=PP_ALIGN.CENTER)
        tx(slide, desc, lx + Inches(0.15), Inches(3.5), Inches(3.6), Inches(1.5),
           18, p["text"], align=PP_ALIGN.CENTER)


def sB04_mikami(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR 01", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "みかみ", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　代表取締役", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    points = [
        ("東大在学中に起業 → AIで組織を最適化", "規模"),
        ("創業5年で年商30億円", "実績"),
        ("AI月商5億・SNS総フォロワー31万人", "影響力"),
        ("受講生6,000名以上「スキルプラス」運営", "信頼"),
    ]
    for i, (point, tag) in enumerate(points):
        ty = Inches(2.65) + i * Inches(1.0)
        if tag:
            box(slide, Inches(0.8), ty, Inches(1.2), Inches(0.7), p["primary"])
            tx(slide, tag, Inches(0.82), ty + Inches(0.2), Inches(1.15), Inches(0.4),
               13, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        tx(slide, point, Inches(2.2), ty + Inches(0.1), Inches(6.6), Inches(0.7),
           20, p["text"])
    add_photo(slide, "mikami-profile.png", Inches(10.1), Inches(0.55), Inches(3.0), Inches(3.0))


def sB05_sato(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR 02", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "佐藤将司", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　AI事業責任者 / AIコンサルタント",
       Inches(0.8), Inches(1.7), Inches(7.5), Inches(0.6), 18, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    facts = [
        ("Claude Codeで100以上のシステムを構築", "実績"),
        ("月間300本の講義動画を自動制作", "自動化"),
        ("法人AIコンサル：月額50万〜200万円で複数社と\n契約中", "単価"),
        ("初月売上4,500万円を突破（4人で達成）", "生産性"),
    ]
    for i, (fact, tag) in enumerate(facts):
        ty = Inches(2.65) + i * Inches(1.0)
        box(slide, Inches(0.8), ty, Inches(1.2), Inches(0.7), p["primary"])
        tx(slide, tag, Inches(0.82), ty + Inches(0.2), Inches(1.15), Inches(0.4),
           13, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        if i == 2:
            tx(slide, fact, Inches(1.833), ty + Inches(0.05), Inches(6.6), Inches(0.9),
               19, p["text"], align=PP_ALIGN.CENTER)
        else:
            tx(slide, fact, Inches(2.2), ty + Inches(0.1), Inches(6.6), Inches(0.7),
               19, p["text"])
    add_photo(slide, "sato-profile.png", Inches(10.1), Inches(0.55), Inches(3.0), Inches(3.0))


def sB06_bootcamp(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "7日間の集中プログラム", Inches(0.8), Inches(0.3),
       Inches(11.7), Inches(0.7), 26, p["text"], bold=True, align=PP_ALIGN.CENTER)
    steps = [
        ("予習期間", "講義動画＋アクションマップで自分のペースで学ぶ", "コミュニティで質問し放題"),
        ("ブートキャンプ週", "月・水・金・土　各2時間（21時〜）の集中作業セッション", "全4回・オンライン"),
        ("補講期間", "ブートキャンプ後1ヶ月間、コミュニティで質問し放題", "アーカイブ動画あり"),
    ]
    for i, (title, desc, sub) in enumerate(steps):
        ty = Inches(1.2) + i * Inches(1.9)
        box(slide, Inches(0.6), ty, Inches(0.7), Inches(1.5), p["primary"])
        tx(slide, str(i + 1), Inches(0.63), ty + Inches(0.4),
           Inches(0.65), Inches(0.8), 30, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        box(slide, Inches(1.5), ty, Inches(11.2), Inches(1.5), p["card"], p["card_bd"], 1)
        tx(slide, title, Inches(1.7), ty + Inches(0.1), Inches(4.0), Inches(0.6),
           18, p["primary"], bold=True)
        tx(slide, desc, Inches(1.7), ty + Inches(0.65), Inches(7.0), Inches(0.6),
           17, p["text"])
        tx(slide, sub, Inches(9.0), ty + Inches(0.35), Inches(3.4), Inches(0.8),
           14, p["muted"], align=PP_ALIGN.CENTER)


def sB07_actionmap(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "7日間で習得できる、実務直結のスキル",
       Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.65), 24, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    themes = [
        ("actionmap-kensakujidoka.png",  "案件探索\n自動化", "外注不要\n営業前に候補が集まる", "時間削減"),
        ("actionmap-slide-seisaku.png",  "スライド\n制作",   "提案資料を\n半自動で作れる",     "工数削減"),
        ("actionmap-system-kaihatsu.png","システム\n開発",   "業務改善ツールを\nAIで内製化",   "外注費削減"),
    ]
    col_w = Inches(4.0)
    img_h = Inches(2.25)
    for i, (fname, theme, value, tag) in enumerate(themes):
        lx = Inches(0.55) + i * Inches(4.27)
        has_img = add_actionmap(slide, fname, lx, Inches(1.0), col_w, img_h)
        if not has_img:
            box(slide, lx, Inches(1.0), col_w, img_h, p["card"], p["card_bd"], 1)
        box(slide, lx, Inches(3.35), col_w, Inches(0.55), p["primary"])
        tx(slide, f"{theme.replace(chr(10), ' ')}  [{tag}]", lx + Inches(0.1), Inches(3.38),
           col_w - Inches(0.2), Inches(0.5), 15, p["text_inv"], bold=True)
        box(slide, lx, Inches(4.0), col_w, Inches(2.2), p["card"], p["card_bd"], 1)
        tx(slide, value, lx + Inches(0.15), Inches(4.15),
           col_w - Inches(0.3), Inches(1.9), 22, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, "全13テーマ収録。社内12種のシステムがそのまま手に入ります。",
       Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), 16, p["primary"],
       bold=True, align=PP_ALIGN.CENTER)


def sB08_proof(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "実績", Inches(0.8), Inches(0.2),
       Inches(11.7), Inches(0.65), 26, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    facts = [
        ("月間300本の動画を", "4人で", "自動制作"),
        ("法人AIコンサル", "月額50〜\n200万円", "で複数社と契約中"),
        ("初月売上", "4,500万円", "を突破"),
    ]
    for i, (pre, num, post) in enumerate(facts):
        lx = Inches(0.5) + i * Inches(4.2)
        box(slide, lx, Inches(1.1), Inches(3.9), Inches(4.5), p["card"], p["card_bd"], 1)
        tx(slide, pre, lx + Inches(0.15), Inches(1.3), Inches(3.6), Inches(0.65),
           16, p["muted"], align=PP_ALIGN.CENTER)
        tx(slide, num, lx + Inches(0.15), Inches(2.0), Inches(3.6), Inches(1.3),
           34, p["primary"], bold=True, align=PP_ALIGN.CENTER)
        post_y = Inches(3.6) if i == 1 else Inches(3.3)
        tx(slide, post, lx + Inches(0.15), post_y, Inches(3.6), Inches(0.65),
           16, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "「机上の空論ではなく、現役で稼ぎ続けている人間が教える」",
       Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), 20, p["text"],
       bold=True, align=PP_ALIGN.CENTER)


def sB09_price(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "投資対効果", Inches(0.8), Inches(0.2),
       Inches(11.7), Inches(0.65), 22, p["muted"], align=PP_ALIGN.CENTER)
    box(slide, Inches(0.6), Inches(1.1), Inches(5.6), Inches(4.8), p["card"], p["card_bd"], 1)
    tx(slide, "法人AIコンサル外注", Inches(0.8), Inches(1.25),
       Inches(5.2), Inches(0.7), 18, p["muted"], bold=True)
    tx(slide, "月50〜200万円", Inches(0.8), Inches(2.0),
       Inches(5.2), Inches(1.2), 44, p["danger"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "× 継続コスト", Inches(0.8), Inches(3.2), Inches(5.2), Inches(0.65),
       18, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "VS", Inches(6.35), Inches(3.0), Inches(0.9), Inches(1.0),
       32, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(7.4), Inches(1.1), Inches(5.6), Inches(4.8), p["card"], p["primary"], 2)
    tx(slide, "このブートキャンプ", Inches(7.6), Inches(1.25),
       Inches(5.2), Inches(0.7), 18, p["primary"], bold=True)
    tx(slide, "198,000円", Inches(7.6), Inches(2.0),
       Inches(5.2), Inches(1.2), 44, p["safe"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "× 買い切り・内製化", Inches(7.6), Inches(3.2),
       Inches(5.2), Inches(0.65), 18, p["primary"], align=PP_ALIGN.CENTER)
    tx(slide, "1日あたり251円。社員1人のAIスキルに投資する価値。",
       Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.65), 20, p["text"],
       align=PP_ALIGN.CENTER)


def sB10_cta(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "自社のAI活用を、7日間で内製化する。",
       Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.9), 28, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.05), p["primary"])
    tx(slide, "本質のClaude Code 完全攻略 7dayブートキャンプ",
       Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.8), 26, p["primary"],
       align=PP_ALIGN.CENTER)
    box(slide, Inches(2.0), Inches(3.3), Inches(9.3), Inches(1.5), p["primary"])
    tx(slide, "お申し込みはこちら →", Inches(2.2), Inches(3.65),
       Inches(8.9), Inches(1.1), 36, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "198,000円（税抜）　月々 7,532円〜",
       Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.65), 20, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "まずは7日間でAI活用の手応えをつかんでください。",
       Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.65), 20, p["text"],
       align=PP_ALIGN.CENTER)


SLIDES_B = [
    sB01_cover, sB02_problem, sB03_solution,
    sB04_mikami, sB05_sato, sB06_bootcamp,
    sB07_actionmap, sB08_proof, sB09_price, sB10_cta,
]


# ===== Pattern C: 初心者安心訴求 =====

def sC01_cover(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "プログラミングは、いりません。",
       Inches(1.0), Inches(1.0), Inches(11.3), Inches(1.5),
       52, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(2.0), Inches(2.7), Inches(9.3), Inches(0.06), p["accent"])
    tx(slide, "日本語で話しかけるだけで動く。",
       Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.85),
       30, p["primary"], align=PP_ALIGN.CENTER)
    tx(slide, "本質のClaude Code 完全攻略 7dayブートキャンプ",
       Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.7), 20, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "アドネス株式会社", Inches(1.0), Inches(6.8),
       Inches(11.3), Inches(0.5), 16, p["muted"], align=PP_ALIGN.RIGHT)


def sC02_problem(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "こんな気持ち、ありませんか？",
       Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.75), 26, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    items = [
        "「毎回同じ作業、なんか時間もったいないな」",
        "「AIって難しそう…自分には無理かも」",
        "「副業してみたいけど、何のスキルを身につければいいか分からない」",
        "「プログラミングは絶対無理。でもAIは気になる」",
    ]
    for i, item in enumerate(items):
        ty = Inches(1.4) + i * Inches(1.3)
        box(slide, Inches(0.6), ty, Inches(12.1), Inches(1.05), p["card"], p["card_bd"], 1)
        tx(slide, item, Inches(0.9), ty + Inches(0.15), Inches(11.5), Inches(0.75),
           22, p["text"])
    tx(slide, "その感覚、全部正解です。だから、このブートキャンプがあります。",
       Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.55), 18, p["primary"],
       bold=True, align=PP_ALIGN.CENTER)


def sC03_solution(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Codeは、", Inches(0.8), Inches(0.5),
       Inches(11.7), Inches(0.85), 34, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "日本語で話しかけるだけで動きます。",
       Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.3),
       44, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(2.75), Inches(10.3), Inches(0.05), p["accent"])
    feats = [
        ("プログラミング不要", "日本語で指示するだけ"),
        ("副業に直結", "案件探索・LP制作・資料作成を自動化"),
        ("未経験OK", "1期生は7日でシステムを完成させた"),
    ]
    for i, (title, desc) in enumerate(feats):
        ty = Inches(3.0) + i * Inches(1.3)
        box(slide, Inches(0.6), ty, Inches(0.06), Inches(0.9), p["primary"])
        tx(slide, title, Inches(0.9), ty, Inches(4.0), Inches(0.55),
           20, p["primary"], bold=True)
        tx(slide, desc, Inches(0.9), ty + Inches(0.5), Inches(11.5), Inches(0.6),
           19, p["text"])


def sC04_mikami(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR 01", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "みかみ", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　代表取締役", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    points = [
        "受講生6,000名以上「スキルプラス」を運営",
        "AIで月商5億円を達成（本物の実績）",
        "東大在学中に起業・創業5年で年商30億円",
        "「現役で稼いでいる人間が教える」を徹底",
    ]
    for i, point in enumerate(points):
        ty = Inches(2.65) + i * Inches(1.0)
        box(slide, Inches(0.8), ty + Inches(0.25), Inches(0.5), Inches(0.5), p["primary"])
        # 上から3項目は13px右にずらし、4番目（「から始まる行）はそのまま
        tx_x = Inches(1.635) if i < 3 else Inches(1.5)
        tx(slide, point, tx_x, ty + Inches(0.25), Inches(7.3), Inches(0.5),
           20, p["text"])
    add_photo(slide, "mikami-profile.png", Inches(10.1), Inches(0.55), Inches(3.0), Inches(3.0))


def sC05_sato(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR 02", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "佐藤将司", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　AI事業責任者", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    tx(slide, "プログラミング経験ゼロからスタートした人が教えます。",
       Inches(0.8), Inches(2.7), Inches(8.0), Inches(0.75), 20, p["primary"], bold=True)
    story = [
        "元・公務員（月給12万円）／プログラミング経験ゼロ",
        "↓　Claude Codeで100以上のシステムを構築",
        "↓　未経験からAIで稼げるようになった",
    ]
    colors = [p["muted"], p["text"], p["primary"]]
    for i, (line, color) in enumerate(zip(story, colors)):
        tx(slide, line, Inches(1.0), Inches(3.5) + i * Inches(1.0),
           Inches(8.0), Inches(0.85), 21, color, bold=(i == 2))
    tx(slide, "→ 「自分でもできるかも」と思えるはずです",
       Inches(0.8), Inches(6.9), Inches(8.2), Inches(0.4), 15, p["muted"],
       align=PP_ALIGN.RIGHT)
    add_photo(slide, "sato-profile.png", Inches(10.1), Inches(0.55), Inches(3.0), Inches(3.0))


def sC06_support(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "一人にしません。", Inches(0.8), Inches(0.3),
       Inches(11.7), Inches(1.0), 48, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    supports = [
        ("予習期間", "講義動画を自分のペースで\n視聴\nコミュニティで質問し放題"),
        ("ブートキャンプ週", "週4回・各2時間のオンラインセッション\n講師と一緒に手を動かす"),
        ("補講期間", "終了後1ヶ月間サポート継続\nアーカイブ動画でいつでも復習"),
    ]
    for i, (title, desc) in enumerate(supports):
        lx = Inches(0.5) + i * Inches(4.2)
        box(slide, lx, Inches(1.6), Inches(3.9), Inches(4.5), p["card"], p["card_bd"], 2)
        box(slide, lx, Inches(1.6), Inches(3.9), Inches(0.65), p["primary"])
        tx(slide, title, lx + Inches(0.1), Inches(1.62), Inches(3.7), Inches(0.55),
           18, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        tx(slide, desc, lx + Inches(0.15), Inches(2.4), Inches(3.6), Inches(2.5),
           19, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, "「分からなくても大丈夫」という環境を作っています。",
       Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.65), 20, p["text"],
       bold=True, align=PP_ALIGN.CENTER)


def sC07_actionmap(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "7日後、こんなことができるようになります",
       Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.65), 24, p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    themes = [
        ("actionmap-kensakujidoka.png",  "案件探索\n自動化", "仕事探しに\n時間をかけなくていい"),
        ("actionmap-slide-seisaku.png",  "スライド\n制作",   "毎回ゼロから\n作る必要がなくなる"),
        ("actionmap-system-kaihatsu.png","システム\n開発",   "業務改善ツールを\n自分で作れる"),
    ]
    col_w = Inches(4.0)
    img_h = Inches(2.25)
    for i, (fname, theme, value) in enumerate(themes):
        lx = Inches(0.55) + i * Inches(4.27)
        has_img = add_actionmap(slide, fname, lx, Inches(1.0), col_w, img_h)
        if not has_img:
            box(slide, lx, Inches(1.0), col_w, img_h, p["card"], p["card_bd"], 1)
        box(slide, lx, Inches(3.35), col_w, Inches(0.55), p["primary"])
        tx(slide, theme.replace("\n", " "), lx + Inches(0.1), Inches(3.38),
           col_w - Inches(0.2), Inches(0.5), 16, p["text_inv"], bold=True)
        box(slide, lx, Inches(4.0), col_w, Inches(2.2), p["card"], p["card_bd"], 1)
        tx(slide, value, lx + Inches(0.15), Inches(4.15),
           col_w - Inches(0.3), Inches(1.9), 22, p["text"], align=PP_ALIGN.CENTER)


def sC08_proof(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "1期生の声", Inches(0.8), Inches(0.2),
       Inches(11.7), Inches(0.7), 26, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(0.6), Inches(1.05), Inches(11.7), Inches(0.05), p["accent"])
    tx(slide, "完全未経験から", Inches(0.8), Inches(1.3),
       Inches(11.7), Inches(0.7), 26, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "7日でシステムを完成", Inches(0.8), Inches(2.0),
       Inches(11.7), Inches(1.5), 60, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "させた受講生がいます。", Inches(0.8), Inches(3.45),
       Inches(11.7), Inches(0.8), 26, p["muted"], align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "LP制作・SNS自動化・業務効率化など、各自が実際に動くシステムを完成させました。",
       Inches(0.8), Inches(4.65), Inches(11.7), Inches(0.75), 20, p["text"],
       align=PP_ALIGN.CENTER)
    tx(slide, "あなたの「初めての作品」を一緒に作りましょう。",
       Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.65), 20, p["primary"],
       bold=True, align=PP_ALIGN.CENTER)


def sC09_price(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "料金", Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.65),
       22, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "198,000", Inches(0.8), Inches(1.0), Inches(11.7), Inches(2.5),
       110, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "円（税抜）", Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.75),
       30, p["muted"], align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.05), p["primary"])
    tx(slide, "1日あたりたったの", Inches(1.5), Inches(4.4), Inches(5.5), Inches(0.65),
       20, p["muted"], align=PP_ALIGN.RIGHT)
    tx(slide, "251円", Inches(7.1), Inches(4.3), Inches(4.2), Inches(0.85),
       42, p["primary"], bold=True)
    tx(slide, "コンビニのランチより安い。", Inches(0.8), Inches(5.3),
       Inches(11.7), Inches(0.65), 20, p["accent"], align=PP_ALIGN.CENTER)
    tx(slide, "クレジット分割：月々約 7,532円〜", Inches(0.8), Inches(6.1),
       Inches(11.7), Inches(0.65), 20, p["muted"], align=PP_ALIGN.CENTER)


def sC10_cta(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "まずやってみることが、", Inches(0.8), Inches(0.8),
       Inches(11.7), Inches(0.9), 34, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "一番の近道です。", Inches(0.8), Inches(1.65),
       Inches(11.7), Inches(1.2), 52, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "本質のClaude Code 完全攻略 7dayブートキャンプ",
       Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.7), 22, p["text"],
       align=PP_ALIGN.CENTER)
    box(slide, Inches(2.0), Inches(4.1), Inches(9.3), Inches(1.5), p["primary"])
    tx(slide, "申し込んでみる →", Inches(2.0), Inches(4.53),
       Inches(9.3), Inches(0.65), 38, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "198,000円（税抜）　月々 7,532円〜　質問し放題・補講あり",
       Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.65), 18, p["muted"],
       align=PP_ALIGN.CENTER)


def sC09b_reason(slide, p):
    """今やる理由 — 初心者向けの行動喚起"""
    bg(slide, p)
    frame(slide, p)
    tx(slide, "「分からない」を卒業する、", Inches(0.8), Inches(1.2),
       Inches(11.7), Inches(1.2), 44, p["muted"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "最初の7日間です。", Inches(0.8), Inches(2.3),
       Inches(11.7), Inches(1.6), 56, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "AIが得意な人だけが得をする時代が、もうすぐそこまで来ています。",
       Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.7), 20, p["text"],
       align=PP_ALIGN.CENTER)
    tx(slide, "最初の一歩を、一人でなく一緒に踏み出しましょう。",
       Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.7), 20, p["primary"],
       bold=True, align=PP_ALIGN.CENTER)


SLIDES_C = [
    sC01_cover, sC02_problem, sC03_solution,
    sC04_mikami, sC05_sato, sC06_support,
    sC07_actionmap, sC08_proof, sC09_price, sC09b_reason, sC10_cta,
]


# ===== ビルド =====

def build(pattern_key, slides_list, out_dir):
    p = PATTERNS[pattern_key]
    prs = new_prs()
    for fn in slides_list:
        slide = blank(prs)
        fn(slide, p)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = str(out_dir / f"cc_v6_{pattern_key}_{ts}.pptx")
    prs.save(out)
    print(f"✅ Pattern {pattern_key} [{p['name']}] → {out}")
    return out


def main():
    out_dir = Path(__file__).parent / "output"
    out_dir.mkdir(exist_ok=True)
    print("🚀 V6 生成開始（3ターゲット訴求パターン）")
    print("   A: 副業・案件獲得訴求（黒×パープル）")
    print("   B: 法人・業務改善訴求（白×青）")
    print("   C: 初心者安心訴求（白×グリーン）\n")
    build("A", SLIDES_A, out_dir)
    build("B", SLIDES_B, out_dir)
    build("C", SLIDES_C, out_dir)
    print("\n✨ 全完了")


if __name__ == "__main__":
    main()
