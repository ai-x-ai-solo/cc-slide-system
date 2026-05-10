#!/usr/bin/env python3
"""
V5: 3カラーパレット × 16枚固定構成（みかみ・佐藤自己紹介追加）
設計: 1スライド1感情 / 権威→再現性の順 / 思考プロセスを言語化
"""

from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

ASSETS_DIR = Path(__file__).parent / "assets"

# ===== 3カラーパレット =====
PALETTES = {
    "A": {
        "name": "白×淡パープル（クリーン）",
        "bg":      RGBColor(0xFF, 0xFF, 0xFF),
        "primary": RGBColor(0x6A, 0x4C, 0xAF),   # パープル
        "accent":  RGBColor(0x9B, 0x72, 0xCF),   # 淡パープル
        "text":    RGBColor(0x1A, 0x1A, 0x2E),   # 濃紺（テキスト）
        "text_inv":RGBColor(0xFF, 0xFF, 0xFF),
        "muted":   RGBColor(0x88, 0x88, 0x99),
        "card":    RGBColor(0xF2, 0xED, 0xFF),   # 超淡パープル
        "card_bd": RGBColor(0xC0, 0xA0, 0xFF),   # カードボーダー
    },
    "B": {
        "name": "黒×オレンジ（CCブートキャンプ公式）",
        "bg":      RGBColor(0x0D, 0x0D, 0x0D),
        "primary": RGBColor(0xFF, 0x66, 0x00),   # CC オレンジ
        "accent":  RGBColor(0xFF, 0xB3, 0x00),   # ゴールド
        "text":    RGBColor(0xFF, 0xFF, 0xFF),
        "text_inv":RGBColor(0x0D, 0x0D, 0x0D),
        "muted":   RGBColor(0x88, 0x88, 0x88),
        "card":    RGBColor(0x1A, 0x1A, 0x1A),
        "card_bd": RGBColor(0xFF, 0x66, 0x00),
    },
    "C": {
        "name": "濃紺×パープル×ゴールド（プレミアム）",
        "bg":      RGBColor(0x0D, 0x1B, 0x3E),   # アドネス濃紺
        "primary": RGBColor(0x9B, 0x72, 0xCF),   # パープル
        "accent":  RGBColor(0xFF, 0xB3, 0x00),   # ゴールド
        "text":    RGBColor(0xFF, 0xFF, 0xFF),
        "text_inv":RGBColor(0x0D, 0x1B, 0x3E),
        "muted":   RGBColor(0xAA, 0xBB, 0xCC),
        "card":    RGBColor(0x14, 0x2A, 0x5A),
        "card_bd": RGBColor(0x9B, 0x72, 0xCF),
    },
}


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, p):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = p["bg"]


def box(slide, l, t, w, h, color, bd_color=None, bd_w=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if bd_color:
        s.line.color.rgb = bd_color
        if bd_w:
            s.line.width = Pt(bd_w)
    else:
        s.line.fill.background()
    return s


def tx(slide, text, l, top, w, h, size, color,
       bold=False, align=PP_ALIGN.LEFT, wrap=True):
    if not text:
        return
    tb = slide.shapes.add_textbox(l, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    r = para.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Hiragino Kaku Gothic ProN"


def frame(slide, p):
    """上下アクセントライン"""
    box(slide, 0, 0, SLIDE_W, Inches(0.07), p["primary"])
    box(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), p["accent"])


# ===== 16枚のスライド関数 =====

def s01_cover(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "AIを使える人と、", Inches(1.0), Inches(0.9),
       Inches(11.3), Inches(1.0), 44, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "使えない人の差が", Inches(1.0), Inches(1.85),
       Inches(11.3), Inches(1.0), 44, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "もう始まっています。", Inches(1.0), Inches(2.8),
       Inches(11.3), Inches(1.0), 44, p["accent"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(4.0), Inches(4.0), Inches(5.3), Inches(0.05), p["muted"])
    tx(slide, "月3,000円のAIツールで　月10万円案件を狙える時代。",
       Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.75), 24, p["muted"],
       align=PP_ALIGN.CENTER)
    # A・C専用：下部にブートキャンプ名を表示（Bは黒背景なので除外）
    if p["bg"] != RGBColor(0x0D, 0x0D, 0x0D):
        box(slide, 0, Inches(6.6), SLIDE_W, Inches(0.85), p["primary"])
        tx(slide, "本質のClaude Code 完全攻略 7dayブートキャンプ",
           Inches(0.5), Inches(6.68), Inches(12.3), Inches(0.7), 24,
           p["text_inv"], bold=True, align=PP_ALIGN.CENTER)


def s02_problem(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "こんな作業、毎日やっていませんか？", Inches(0.8), Inches(0.3),
       Inches(11.7), Inches(0.8), 30, p["text"] if p["bg"] == RGBColor(0xFF,0xFF,0xFF) else p["text"],
       bold=True, align=PP_ALIGN.CENTER)
    items = ["資料作成", "SNS投稿", "リサーチ", "メール返信", "営業提案", "コピペ作業"]
    for i, item in enumerate(items):
        col, row = i % 2, i // 2
        lx = Inches(1.3) + col * Inches(6.1)
        ty = Inches(1.3) + row * Inches(1.35)
        box(slide, lx, ty, Inches(5.7), Inches(1.05), p["card"], p["card_bd"], 1.5)
        tx(slide, item, lx + Inches(0.25), ty + Inches(0.2),
           Inches(5.2), Inches(0.72), 26, p["text"], bold=True)
    box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.05), p["primary"])
    tx(slide, "その作業、AIで消せます。", Inches(0.8), Inches(5.65),
       Inches(11.7), Inches(0.7), 28, p["primary"], bold=True, align=PP_ALIGN.CENTER)


def s03_contrast(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, 'AIを  "検索"  にしか使えていない人が多い。',
       Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8), 28,
       p["text"], bold=True, align=PP_ALIGN.CENTER)
    # 左：多くの人
    box(slide, Inches(0.4), Inches(1.2), Inches(5.6), Inches(4.5), p["card"], p["muted"], 1)
    tx(slide, "😓  多くの人", Inches(0.6), Inches(1.35), Inches(5.1), Inches(0.6), 18, p["muted"], bold=True)
    tx(slide, "ChatGPTに\n質問する", Inches(0.6), Inches(2.1),
       Inches(5.1), Inches(1.8), 30, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "→ 答えをもらうだけ", Inches(0.6), Inches(4.3),
       Inches(5.1), Inches(0.6), 18, p["muted"], align=PP_ALIGN.CENTER)
    # VS
    tx(slide, "VS", Inches(6.2), Inches(3.5), Inches(1.0), Inches(0.8),
       36, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    # 右：強い人
    box(slide, Inches(7.4), Inches(1.2), Inches(5.5), Inches(4.5), p["card"], p["primary"], 2)
    tx(slide, "⚡  強い人", Inches(7.6), Inches(1.35), Inches(5.0), Inches(0.6), 18, p["primary"], bold=True)
    tx(slide, "AIに\n働かせる", Inches(7.6), Inches(2.1),
       Inches(5.0), Inches(1.8), 30, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "→ 自動で仕事が回る", Inches(7.6), Inches(4.3),
       Inches(5.0), Inches(0.6), 18, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, '本当に強い人は、AIに "働かせて" います。',
       Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), 22,
       p["accent"], bold=True, align=PP_ALIGN.CENTER)


def s04_intro_mikami(slide, p):
    """みかみ自己紹介 → 権威・規模感：「この会社、本当に強い」"""
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "みかみ", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　代表取締役", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    points = [
        ("東大在学中に起業 → AI導入で組織を最適化", "規模"),
        ("AIで月商5億円を達成", "実績"),
        ("スキルプラス受講生 6,000名以上", "信頼"),
        ("「AIを実際の事業で使い続けている側」", ""),
    ]
    for i, (point, tag) in enumerate(points):
        ty = Inches(2.65) + i * Inches(1.0)
        if tag:
            box(slide, Inches(0.8), ty, Inches(1.1), Inches(0.7), p["primary"])
            tx(slide, tag, Inches(0.82), ty + Inches(0.1), Inches(1.05), Inches(0.55),
               13, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
        tx(slide, point, Inches(2.1), ty + Inches(0.1), Inches(6.7), Inches(0.7),
           22, p["text"])
    tx(slide, "→ 「この会社、本物だ」", Inches(0.8), Inches(6.9),
       Inches(8.2), Inches(0.4), 15, p["muted"], align=PP_ALIGN.RIGHT)
    # 写真（右側）
    photo = ASSETS_DIR / "photo_mikami.png"
    if photo.exists():
        slide.shapes.add_picture(str(photo), Inches(9.3), Inches(1.0), Inches(3.5), Inches(3.5))


def s05_intro_sato(slide, p):
    """佐藤自己紹介 → 再現性・親近感：「自分にも出来そう」"""
    bg(slide, p)
    frame(slide, p)
    tx(slide, "INSTRUCTOR", Inches(0.8), Inches(0.2), Inches(4.0), Inches(0.5),
       14, p["muted"], bold=True)
    tx(slide, "佐藤将司", Inches(0.8), Inches(0.65), Inches(7.5), Inches(1.1),
       52, p["primary"], bold=True)
    tx(slide, "アドネス株式会社　AI事業責任者", Inches(0.8), Inches(1.7),
       Inches(7.5), Inches(0.6), 20, p["muted"])
    box(slide, Inches(0.8), Inches(2.45), Inches(8.2), Inches(0.05), p["accent"])
    story = [
        "元・公務員（月給12万円）／プログラミング経験ゼロ",
        "↓　Claude Codeで100以上のシステムを構築",
        "↓　営業初月　1,000万円案件を獲得",
        "↓　AIで月4,500万円規模の事業を運営",
    ]
    colors = [p["muted"], p["text"], p["text"], p["primary"]]
    for i, (line, color) in enumerate(zip(story, colors)):
        bold = i == 3
        tx(slide, line, Inches(1.0), Inches(2.7) + i * Inches(1.0),
           Inches(8.0), Inches(0.85), 22 if i < 3 else 24, color, bold=bold)
    tx(slide, "→ 「未経験でも、自分にも出来るかも」", Inches(0.8), Inches(6.9),
       Inches(8.2), Inches(0.4), 15, p["muted"], align=PP_ALIGN.RIGHT)
    # 写真（右側）
    photo = ASSETS_DIR / "photo_sato.png"
    if photo.exists():
        slide.shapes.add_picture(str(photo), Inches(9.3), Inches(1.0), Inches(3.5), Inches(3.5))


def s06_bootcamp(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Code ブートキャンプ", Inches(1.0), Inches(0.3),
       Inches(11.3), Inches(0.65), 26, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "7", Inches(1.0), Inches(0.95), Inches(11.3), Inches(3.5),
       160, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "DAYS", Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.1),
       64, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(0.05), p["accent"])
    points = ["未経験OK", "日本語で操作", "実践型", "その場で完成"]
    for i, pt in enumerate(points):
        tx(slide, f"✓ {pt}", Inches(0.8) + i * Inches(3.0), Inches(5.4),
           Inches(2.8), Inches(0.7), 20, p["text"], align=PP_ALIGN.CENTER)


def s07_curriculum(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "実際に学べる内容", Inches(0.8), Inches(0.25),
       Inches(11.7), Inches(0.7), 28, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "AIエージェント ×", Inches(0.8), Inches(0.9),
       Inches(11.7), Inches(0.65), 24, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    left_items =  ["LP制作 自動化", "案件探索 自動化", "X投稿 自動化", "LINE Bot 構築"]
    right_items = ["動画編集 自動化", "スライド制作", "システム開発", "SEO記事生成"]
    for i, item in enumerate(left_items):
        ty = Inches(1.7) + i * Inches(1.15)
        box(slide, Inches(0.4), ty, Inches(5.9), Inches(0.95), p["card"], p["card_bd"], 1.5)
        tx(slide, item, Inches(0.7), ty + Inches(0.15), Inches(5.3), Inches(0.75), 24, p["text"], bold=True)
    for i, item in enumerate(right_items):
        ty = Inches(1.7) + i * Inches(1.15)
        box(slide, Inches(7.0), ty, Inches(5.9), Inches(0.95), p["card"], p["card_bd"], 1.5)
        tx(slide, item, Inches(7.3), ty + Inches(0.15), Inches(5.3), Inches(0.75), 24, p["text"], bold=True)


def s08_proof(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "机上の空論ではありません。", Inches(0.8), Inches(0.3),
       Inches(11.7), Inches(0.7), 26, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "月4,500万円", Inches(0.3), Inches(0.9),
       Inches(12.7), Inches(2.4), 80, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "を実際に生み出しているチームが教えます。（4人で達成）",
       Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.7), 22, p["text"], align=PP_ALIGN.CENTER)
    box(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.05), p["accent"])
    proofs = ["プログラミング未経験から100以上のシステムを構築", "AIで講義300本 / 月 自動制作", "営業初月 1,000万円案件を獲得"]
    for i, proof in enumerate(proofs):
        tx(slide, f"▶  {proof}", Inches(2.0), Inches(4.2) + i * Inches(0.7),
           Inches(9.5), Inches(0.6), 22, p["muted"])


def s09_before_after(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "未経験でも大丈夫？", Inches(0.8), Inches(0.25),
       Inches(11.7), Inches(0.75), 32, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "大丈夫です。", Inches(0.8), Inches(0.95),
       Inches(11.7), Inches(0.95), 48, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(0.4), Inches(2.1), Inches(5.6), Inches(3.6), p["card"], p["muted"], 1)
    tx(slide, "Before", Inches(0.6), Inches(2.2), Inches(5.1), Inches(0.55), 17, p["muted"], bold=True)
    tx(slide, "「何から始めれば\nいいかわからない…」",
       Inches(0.6), Inches(2.85), Inches(5.1), Inches(1.9), 26, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "7日後 →", Inches(6.1), Inches(3.65), Inches(1.3), Inches(0.7),
       20, p["accent"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(7.4), Inches(2.1), Inches(5.5), Inches(3.6), p["card"], p["primary"], 2)
    tx(slide, "After", Inches(7.6), Inches(2.2), Inches(5.0), Inches(0.55), 17, p["primary"], bold=True)
    tx(slide, "「自分のシステムが\n動いてる！」",
       Inches(7.6), Inches(2.85), Inches(5.0), Inches(1.9), 28, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "1期生：完全未経験 → 7日でシステム完成",
       Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), 22, p["accent"], bold=True, align=PP_ALIGN.CENTER)


def s10_value(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "あなたが受け取るもの", Inches(0.8), Inches(0.25),
       Inches(11.7), Inches(0.7), 26, p["text"], bold=True, align=PP_ALIGN.CENTER)
    left_items =  ["講義13テーマ（毎週増加中）", "専用コミュニティ", "ブートキャンプ4回（各2時間）"]
    right_items = ["社内システム12種", "1on1コンサル（2時間）", "補講サポート1ヶ月"]
    for i, item in enumerate(left_items):
        tx(slide, f"✓  {item}", Inches(0.8), Inches(1.1) + i * Inches(0.75),
           Inches(5.8), Inches(0.65), 20, p["muted"])
    for i, item in enumerate(right_items):
        tx(slide, f"✓  {item}", Inches(7.0), Inches(1.1) + i * Inches(0.75),
           Inches(5.8), Inches(0.65), 20, p["muted"])
    box(slide, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.05), p["primary"])
    tx(slide, "合計", Inches(0.8), Inches(3.65), Inches(11.7), Inches(0.6),
       24, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "¥2,485,000相当", Inches(0.3), Inches(4.25),
       Inches(12.7), Inches(2.0), 72, p["primary"], bold=True, align=PP_ALIGN.CENTER)


def s11_price(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "しかし今回、2期生限定で。", Inches(0.8), Inches(0.4),
       Inches(11.7), Inches(0.7), 24, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "クレジット分割なら", Inches(0.8), Inches(1.25),
       Inches(11.7), Inches(0.65), 24, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, "月々 約7,532円〜", Inches(0.3), Inches(1.9),
       Inches(12.7), Inches(2.2), 72, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(3.0), Inches(4.2), Inches(7.3), Inches(0.05), p["muted"])
    tx(slide, "1日あたり 約251円", Inches(0.8), Inches(4.4),
       Inches(11.7), Inches(0.75), 32, p["text"], align=PP_ALIGN.CENTER)
    tx(slide, "コーヒー1杯以下で、月10万円案件を狙えるスキルが身につく。",
       Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.65), 20, p["muted"], align=PP_ALIGN.CENTER)


def s12_roi(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "198,000円。", Inches(0.8), Inches(0.35),
       Inches(11.7), Inches(1.1), 52, p["text"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "高いと思いましたか？", Inches(0.8), Inches(1.4),
       Inches(11.7), Inches(0.65), 28, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "でももし、", Inches(0.8), Inches(2.3),
       Inches(11.7), Inches(0.55), 22, p["text"], align=PP_ALIGN.CENTER)
    ifs = ["月10万円案件を2件取ったら", "AIで業務を週10時間削減できたら", "自動化スキルで副業が軌道に乗ったら"]
    for i, line in enumerate(ifs):
        tx(slide, f"▶  {line}", Inches(2.5), Inches(2.95) + i * Inches(0.75),
           Inches(8.5), Inches(0.65), 22, p["text"])
    box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.05), p["primary"])
    tx(slide, "回収できる金額は、それ以上かもしれません。",
       Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7), 26,
       p["primary"], bold=True, align=PP_ALIGN.CENTER)


def s13_urgency(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "Claude Codeを仕事に使いこなしている日本人は",
       Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.65), 24, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "0.1%", Inches(0.3), Inches(0.85),
       Inches(12.7), Inches(3.8), 148, p["primary"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "未満", Inches(0.3), Inches(4.5),
       Inches(12.7), Inches(0.9), 48, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.05), p["accent"])
    tx(slide, "今はまだ、「早い者勝ち」の段階です。",
       Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.65),
       26, p["accent"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "先に動いた人だけが、圧倒的に有利になる。",
       Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.65), 20, p["muted"],
       align=PP_ALIGN.CENTER)


def s14_future(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "もし1年後、AIを使える人だけが生産性を10倍にしていたら？",
       Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.75), 22, p["muted"],
       align=PP_ALIGN.CENTER)
    box(slide, Inches(0.4), Inches(1.2), Inches(5.6), Inches(4.6), p["card"], p["muted"], 1)
    tx(slide, "そのままの自分", Inches(0.6), Inches(1.35), Inches(5.1), Inches(0.6), 19, p["muted"], bold=True)
    for i, f in enumerate(["今と同じ作業を繰り返す", "AIが得意な人に追い越される", "差が広がるのを見ている"]):
        tx(slide, f"✗  {f}", Inches(0.7), Inches(2.1) + i * Inches(0.9),
           Inches(5.1), Inches(0.8), 20, p["muted"])
    tx(slide, "あなたは\nどちら側に\nいますか？", Inches(6.1), Inches(2.8),
       Inches(1.3), Inches(2.5), 18, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(7.4), Inches(1.2), Inches(5.5), Inches(4.6), p["card"], p["primary"], 2)
    tx(slide, "AIを使いこなす自分", Inches(7.6), Inches(1.35), Inches(5.0), Inches(0.6), 19, p["primary"], bold=True)
    for i, f in enumerate(["自動化で時間が生まれる", "月10万円以上の案件を受注", "自由な働き方を実現"]):
        tx(slide, f"✓  {f}", Inches(7.7), Inches(2.1) + i * Inches(0.9),
           Inches(5.0), Inches(0.8), 20, p["text"])
    tx(slide, "決断するなら、「やってみたい」と思っている今です。",
       Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.65), 20, p["muted"],
       align=PP_ALIGN.CENTER)


def s15_cta(slide, p):
    bg(slide, p)
    frame(slide, p)
    tx(slide, "本質のClaude Code 完全攻略", Inches(0.8), Inches(0.7),
       Inches(11.7), Inches(0.7), 26, p["muted"], align=PP_ALIGN.CENTER)
    tx(slide, "7dayブートキャンプ", Inches(0.8), Inches(1.35),
       Inches(11.7), Inches(1.0), 44, p["text"], bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.05), p["primary"])
    tx(slide, "未経験から、AIで仕事を作れる人へ。",
       Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.65), 24, p["muted"],
       align=PP_ALIGN.CENTER)
    box(slide, Inches(2.3), Inches(3.7), Inches(8.7), Inches(1.4), p["primary"])
    tx(slide, "本日参加受付中 →", Inches(2.5), Inches(3.9),
       Inches(8.3), Inches(1.0), 36, p["text_inv"], bold=True, align=PP_ALIGN.CENTER)
    tx(slide, "198,000円（税抜）　クレジット分割：月々約7,532円〜",
       Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.6), 20, p["muted"],
       align=PP_ALIGN.CENTER)
    tx(slide, "※ 2期生限定価格", Inches(0.8), Inches(6.15),
       Inches(11.7), Inches(0.5), 16, p["muted"], align=PP_ALIGN.CENTER)


SLIDES = [
    s01_cover, s02_problem, s03_contrast,
    s04_intro_mikami, s05_intro_sato,      # ← 自己紹介（Slide 4・5）
    s06_bootcamp, s07_curriculum, s08_proof,
    s09_before_after, s10_value, s11_price,
    s12_roi, s13_urgency, s14_future, s15_cta,
]


def build(palette_key: str, out_dir: Path):
    p = PALETTES[palette_key]
    prs = new_prs()
    for fn in SLIDES:
        slide = blank(prs)
        fn(slide, p)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = str(out_dir / f"cc_v5_{palette_key}_{ts}.pptx")
    prs.save(out)
    print(f"✅ Palette {palette_key} [{p['name']}] → {out}")
    return out


def main():
    out_dir = Path(__file__).parent / "output"
    out_dir.mkdir(exist_ok=True)
    print("🚀 V5 生成開始（3パレット × 15枚）")
    for key in ["A", "B", "C"]:
        build(key, out_dir)
    print("\n✨ 全完了")


if __name__ == "__main__":
    main()
