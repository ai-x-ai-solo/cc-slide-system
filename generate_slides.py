#!/usr/bin/env python3
"""
Claude Code ブートキャンプ セールススライド自動生成システム v2
設計原則: 1スライド1メッセージ / 余白を怖がらない / 数字を主役に / 感情を1個だけ動かす
使い方: python3 generate_slides.py --style v1 --emphasis price
"""

import argparse
import json
import yaml
from pathlib import Path
from datetime import datetime
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ===== スタイル定義 =====
STYLES = {
    "v1": {
        "name": "CCオリジナル（白・青）",
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "primary": RGBColor(0x1A, 0x73, 0xC8),   # 青
        "accent": RGBColor(0xFF, 0x66, 0x00),      # オレンジ
        "text": RGBColor(0x1A, 0x1A, 0x1A),
        "text_inv": RGBColor(0xFF, 0xFF, 0xFF),
        "muted": RGBColor(0x77, 0x77, 0x77),
        "card": RGBColor(0xF0, 0xF4, 0xFF),
    },
    "v2": {
        "name": "アドネス公式（濃紺・オレンジ）",
        "bg": RGBColor(0x0D, 0x1B, 0x3E),          # 濃紺
        "primary": RGBColor(0xFF, 0x66, 0x00),      # オレンジ
        "accent": RGBColor(0xFF, 0xB3, 0x00),       # ゴールド
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "text_inv": RGBColor(0x0D, 0x1B, 0x3E),
        "muted": RGBColor(0xAA, 0xBB, 0xCC),
        "card": RGBColor(0x14, 0x2A, 0x5A),
    },
    "v3": {
        "name": "ミニマル（黒・白・金）",
        "bg": RGBColor(0x0A, 0x0A, 0x0A),          # ほぼ黒
        "primary": RGBColor(0xFF, 0xFF, 0xFF),      # 白
        "accent": RGBColor(0xD4, 0xAF, 0x37),       # ゴールド
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "text_inv": RGBColor(0x0A, 0x0A, 0x0A),
        "muted": RGBColor(0x88, 0x88, 0x88),
        "card": RGBColor(0x1A, 0x1A, 0x1A),
    },
}


def load_service_info(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def generate_slide_content(service: dict, emphasis: str, style_key: str) -> list[dict]:
    """1スライド1メッセージの原則でClaude APIがコンテンツを生成"""
    client = anthropic.Anthropic()

    emphasis_map = {
        "price":    "「月3,000円のClaude Codeで月10万円案件が取れる」ROIと価格の安さ",
        "result":   "「4人で初月4,500万円」という圧倒的な実績と数字",
        "beginner": "「未経験でも7日でシステム完成」という安心感",
        "urgency":  "「日本人の0.1%しか使えていない、今が早い者勝ち」という緊急性",
        "value":    "「¥2,485,000相当を198,000円で手に入れる」圧倒的コスパ",
    }
    emphasis_text = emphasis_map.get(emphasis, emphasis_map["result"])

    style_tone = {
        "v1": "清潔感のある青・白。信頼と誠実さ。",
        "v2": "アドネスの濃紺にオレンジアクセント。力強さと実績。",
        "v3": "黒背景にゴールド。プレミアム感と希少性。",
    }

    prompt = f"""
あなたはトップクラスのセールスコピーライター＆スライドデザイナーです。

【絶対に守る設計原則】
1. 1スライド1メッセージ：各スライドで伝えることは1つだけ
2. 数字を主役に：抽象的な言葉より具体的な数字が強い
3. 余白の設計：情報は少ないほど伝わる
4. 感情を1個動かす：各スライドで起こす感情を1つ決める
5. 引き算：「何を削るか」で良し悪しが決まる

【訴求メイン軸】
{emphasis_text}

【スタイルトーン】
{style_tone.get(style_key, '')}

【サービス情報（全情報。どう使うかはあなたが判断）】
{json.dumps(service, ensure_ascii=False, indent=2)}

【出力形式】
以下のJSON配列のみ出力。説明文不要。12〜15枚。

[
  {{
    "slide_type": "cover",
    "headline": "5〜10文字の強烈なキャッチコピー",
    "big_number": "このスライドの主役になる数字（任意）",
    "subtext": "10〜20文字のサポートコピー（任意）",
    "emotion": "このスライドで起こす感情（例：驚き・安心・焦り）"
  }},
  {{
    "slide_type": "number",
    "headline": "数字の意味（10文字以内）",
    "big_number": "インパクトある数字（例：月4,500万円）",
    "subtext": "補足（20文字以内）",
    "emotion": "驚き"
  }},
  {{
    "slide_type": "message",
    "headline": "このスライドで伝える1つのメッセージ（15文字以内）",
    "body": "1〜2文の補足（40文字以内）",
    "emotion": "共感・焦り・安心など"
  }},
  {{
    "slide_type": "before_after",
    "headline": "変化のタイトル（10文字以内）",
    "before": ["今の状態1", "今の状態2", "今の状態3"],
    "after": ["変化後1", "変化後2", "変化後3"],
    "emotion": "希望"
  }},
  {{
    "slide_type": "cta",
    "headline": "行動を促す一言（10文字以内）",
    "subtext": "申し込みへの背中を押すコピー（30文字以内）",
    "button_text": "ボタンに書くテキスト",
    "emotion": "決断"
  }}
]

slide_typeの種類：
- cover：表紙
- number：数字を主役にしたインパクトスライド
- message：1メッセージを大きく見せるスライド
- before_after：BEFORE→AFTER対比スライド
- proof：実績・証拠スライド
- price：価格スライド
- cta：クロージング

流れ：cover → number(実績) → message(課題) → message(解決策) → number(価値) → before_after → message(安心) → price → message(緊急性) → cta

必ずJSON配列のみ返すこと。
"""

    message = client.messages.create(
        model="claude-opus-4-7",
        max_tokens=4000,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = message.content[0].text.strip()
    if "```" in raw:
        parts = raw.split("```")
        for part in parts:
            part = part.strip()
            if part.startswith("json"):
                part = part[4:].strip()
            if part.startswith("["):
                raw = part
                break
    raw = raw.strip()
    start = raw.find("[")
    end = raw.rfind("]")
    if start != -1 and end != -1:
        raw = raw[start:end + 1]

    return json.loads(raw)


# ===== レンダラー =====

def bg(slide, style):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = style["bg"]


def txt(slide, text, left, top, width, height, size, color,
        bold=False, align=PP_ALIGN.LEFT, wrap=True):
    if not text:
        return
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Hiragino Kaku Gothic ProN"


def rect(slide, left, top, width, height, fill_color, line_color=None, line_w=None):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        if line_w:
            s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def render_cover(slide, data, style):
    bg(slide, style)
    # アクセントライン（上部）
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    # ヘッドライン（大）
    txt(slide, data.get("headline", ""), Inches(1.2), Inches(1.5),
        Inches(10.9), Inches(3.0), 54, style["primary"],
        bold=True, align=PP_ALIGN.CENTER)
    # big_number
    if data.get("big_number"):
        txt(slide, data["big_number"], Inches(1.2), Inches(4.3),
            Inches(10.9), Inches(1.2), 44, style["accent"],
            bold=True, align=PP_ALIGN.CENTER)
    # subtext
    if data.get("subtext"):
        txt(slide, data["subtext"], Inches(1.5), Inches(5.6),
            Inches(10.3), Inches(0.8), 22, style["muted"],
            align=PP_ALIGN.CENTER)
    # アクセントライン（下部）
    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


def render_number(slide, data, style):
    bg(slide, style)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    # ラベル（小）
    txt(slide, data.get("headline", ""), Inches(1.0), Inches(0.8),
        Inches(11.3), Inches(0.6), 20, style["muted"],
        align=PP_ALIGN.CENTER)
    # BIG NUMBER（超大）
    txt(slide, data.get("big_number", ""), Inches(0.5), Inches(1.5),
        Inches(12.3), Inches(3.5), 96, style["accent"],
        bold=True, align=PP_ALIGN.CENTER)
    # subtext
    if data.get("subtext"):
        txt(slide, data["subtext"], Inches(1.5), Inches(5.2),
            Inches(10.3), Inches(0.8), 24, style["text"],
            align=PP_ALIGN.CENTER)
    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


def render_message(slide, data, style):
    bg(slide, style)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    # メインメッセージ（大）
    txt(slide, data.get("headline", ""), Inches(1.0), Inches(1.8),
        Inches(11.3), Inches(2.8), 48, style["primary"],
        bold=True, align=PP_ALIGN.CENTER)
    # ボディ
    if data.get("body"):
        txt(slide, data["body"], Inches(1.5), Inches(4.8),
            Inches(10.3), Inches(1.2), 22, style["muted"],
            align=PP_ALIGN.CENTER)
    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


def render_before_after(slide, data, style):
    bg(slide, style)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    txt(slide, data.get("headline", ""), Inches(0.5), Inches(0.2),
        Inches(12.3), Inches(0.6), 22, style["muted"],
        align=PP_ALIGN.CENTER)

    # BEFORE カラム
    rect(slide, Inches(0.4), Inches(1.0), Inches(5.5), Inches(5.8),
         style["card"], style["muted"], 1)
    txt(slide, "BEFORE", Inches(0.5), Inches(1.1), Inches(5.3), Inches(0.6),
        18, style["muted"], bold=True, align=PP_ALIGN.CENTER)
    for i, line in enumerate(data.get("before", [])[:4]):
        txt(slide, f"✗  {line}", Inches(0.7), Inches(1.9) + i * Inches(1.1),
            Inches(5.0), Inches(1.0), 20, style["muted"])

    # 矢印
    txt(slide, "→", Inches(6.1), Inches(3.5), Inches(1.1), Inches(0.8),
        40, style["accent"], bold=True, align=PP_ALIGN.CENTER)

    # AFTER カラム
    rect(slide, Inches(7.4), Inches(1.0), Inches(5.5), Inches(5.8),
         style["card"], style["primary"], 2)
    txt(slide, "AFTER", Inches(7.5), Inches(1.1), Inches(5.3), Inches(0.6),
        18, style["primary"], bold=True, align=PP_ALIGN.CENTER)
    for i, line in enumerate(data.get("after", [])[:4]):
        txt(slide, f"✓  {line}", Inches(7.6), Inches(1.9) + i * Inches(1.1),
            Inches(5.0), Inches(1.0), 20, style["text"])

    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


def render_proof(slide, data, style):
    bg(slide, style)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    txt(slide, data.get("headline", ""), Inches(0.5), Inches(0.2),
        Inches(12.3), Inches(0.7), 26, style["primary"],
        bold=True, align=PP_ALIGN.CENTER)
    body = data.get("body", []) if isinstance(data.get("body"), list) else []
    for i, line in enumerate(body[:5]):
        top = Inches(1.2) + i * Inches(1.0)
        rect(slide, Inches(0.5), top, Inches(12.3), Inches(0.85),
             style["card"])
        txt(slide, line, Inches(0.8), top + Inches(0.1), Inches(11.8),
            Inches(0.7), 22, style["text"])
    if data.get("subtext"):
        txt(slide, data["subtext"], Inches(1.0), Inches(6.6),
            Inches(11.3), Inches(0.6), 18, style["accent"],
            bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


def render_price(slide, data, style):
    bg(slide, style)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    txt(slide, data.get("headline", ""), Inches(0.5), Inches(0.2),
        Inches(12.3), Inches(0.6), 24, style["muted"],
        align=PP_ALIGN.CENTER)
    # BIG PRICE
    if data.get("big_number"):
        txt(slide, data["big_number"], Inches(0.5), Inches(1.1),
            Inches(12.3), Inches(2.2), 72, style["accent"],
            bold=True, align=PP_ALIGN.CENTER)
    body = data.get("body", []) if isinstance(data.get("body"), list) else []
    for i, line in enumerate(body[:4]):
        txt(slide, line, Inches(1.5), Inches(3.5) + i * Inches(0.75),
            Inches(10.3), Inches(0.65), 22, style["text"],
            align=PP_ALIGN.CENTER)
    if data.get("subtext"):
        rect(slide, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.75),
             style["primary"])
        txt(slide, data["subtext"], Inches(2.2), Inches(6.35),
            Inches(9.0), Inches(0.65), 20, style["text_inv"],
            bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


def render_cta(slide, data, style):
    bg(slide, style)
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), style["primary"])
    txt(slide, data.get("headline", ""), Inches(0.8), Inches(1.2),
        Inches(11.7), Inches(2.5), 56, style["primary"],
        bold=True, align=PP_ALIGN.CENTER)
    if data.get("subtext"):
        txt(slide, data["subtext"], Inches(1.2), Inches(3.9),
            Inches(10.9), Inches(1.0), 24, style["muted"],
            align=PP_ALIGN.CENTER)
    # CTAボタン
    btn_text = data.get("button_text", "今すぐ申し込む →")
    rect(slide, Inches(3.0), Inches(5.3), Inches(7.3), Inches(1.3),
         style["primary"])
    txt(slide, btn_text, Inches(3.2), Inches(5.45),
        Inches(6.9), Inches(1.0), 28, style["text_inv"],
        bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.08), style["accent"])


RENDER_MAP = {
    "cover": render_cover,
    "number": render_number,
    "message": render_message,
    "before_after": render_before_after,
    "proof": render_proof,
    "price": render_price,
    "cta": render_cta,
}


def build_pptx(slides_data: list[dict], style_key: str, output_path: str):
    style = STYLES[style_key]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for data in slides_data:
        slide = prs.slides.add_slide(blank)
        slide_type = data.get("slide_type", "message")
        renderer = RENDER_MAP.get(slide_type, render_message)
        renderer(slide, data, style)

    prs.save(output_path)
    print(f"✅ {output_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--style", choices=["v1", "v2", "v3"], default="v1")
    parser.add_argument("--emphasis",
                        choices=["price", "result", "beginner", "urgency", "value"],
                        default="result")
    parser.add_argument("--service", default="service_info.yaml")
    args = parser.parse_args()

    base = Path(__file__).parent
    service = load_service_info(base / args.service)

    print(f"🚀 [{args.style}] 生成開始 emphasis={args.emphasis}")
    print("📝 Claude APIでコンテンツ生成中...")
    slides_data = generate_slide_content(service, args.emphasis, args.style)
    print(f"   → {len(slides_data)}枚")

    out_dir = base / "output"
    out_dir.mkdir(exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = str(out_dir / f"cc_{args.style}_{args.emphasis}_{ts}.pptx")

    print("🎨 PPTX構築中...")
    build_pptx(slides_data, args.style, out)
    print(f"\n   スタイル: {STYLES[args.style]['name']}")
    print(f"   枚数: {len(slides_data)}枚")
    print(f"   ファイル: {out}")


if __name__ == "__main__":
    main()
